from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import src.addphoto as addphoto
import re
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import random
import re
import requests
import io

def add_slide_numbers(ppt_path):
    # Load the existing PowerPoint presentation
    prs = Presentation(ppt_path)

    for i, slide in enumerate(prs.slides, start=1):
        if i==1 or i==2:
            continue
        left = Inches(6.43)  # Adjust the position as needed
        top = Inches(6.74)   # Adjust the position as needed
        width = Inches(1)
        height = Inches(0.2)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = f"{i}"
        p.font.bold = True
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.RIGHT  # Adjust alignment as needed

    # Save the presentation with slide numbers
    prs.save(ppt_path)

def set_font(ppt_path, selected_font):
    prs = Presentation(ppt_path)
    # Iterate through all slides
    for slide in prs.slides:
        # Iterate through all shapes on the slide
        for shape in slide.shapes:
            # Check if the shape has text
            if shape.has_text_frame:
                # Iterate through all paragraphs in the text frame
                for paragraph in shape.text_frame.paragraphs:
                    # Iterate through all runs in the paragraph
                    for run in paragraph.runs:
                        # Change font characteristics
                        run.font.name = selected_font  # Change the font type
                        # run.font.size = Pt(24)  # Change the font size
                        # run.font.bold = True  # Make the text bold
                        # run.font.color.rgb = RGBColor(0x42, 0x87, 0xF5)  # Change the font color

    # Save your presentation
    prs.save(ppt_path)

def create_ppt(text_file, design_number, ppt_name, author):
    prs = Presentation(f"Designs/Design-{design_number}.pptx")
    slide_master = prs.slide_masters[0]
    slide_count = 0
    header = ""
    content = ""
    last_slide_layout_index = -1
    firsttime = True
    with open(text_file, 'r', encoding='utf-8') as f:
        for line_num, line in enumerate(f):
                            
            if line.startswith('#Title:'):
                header = line.replace('#Title:', '').strip()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                title.text = header
                subtitle = slide.placeholders[1]
                subtitle.text = author
                body_shape = slide.shapes.placeholders[1]
                continue
            elif line.startswith('#Slide:'):
                if slide_count > 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
                    title = slide.shapes.title
                    title.text = header
                    body_shape = slide.shapes.placeholders[slide_placeholder_index]
                    tf = body_shape.text_frame
                    tf.text = content
                    if slide_layout_index == 8:
                        placeholder = slide.placeholders[1]
                        query = f'{ppt_name} {header}'
                        picture = placeholder.insert_picture("images/"+addphoto.get_images(query,1)[0])
                        addphoto.empty_images()
                        # image_url = addphoto.search_pexels_images(ppt_name)
                        # if image_url is not None:
                        #     # download the image
                        #     image_data = requests.get(image_url).content
                        #     # load image into BytesIO object
                        #     image_stream = io.BytesIO(image_data)
                        #     picture = placeholder.insert_picture(image_stream)
                    # if slide_layout_index == 7:
                    #     placeholder = slide.placeholders[]
                    #     query = f'{ppt_name} {header}'
                    #     picture = placeholder.insert_picture("images/"+addphoto.get_images(query,1)[0])
                    #     addphoto.empty_images()
                    

                content = "" 
                slide_count += 1
                slide_layout_index = last_slide_layout_index
                layout_indices = [1, 7, 8] 
                while slide_layout_index == last_slide_layout_index:
                    if firsttime == True:
                        slide_layout_index = 1
                        slide_placeholder_index = 1
                        firsttime = False
                        break
                    slide_layout_index = random.choice(layout_indices) # Select random slide index
                    if slide_layout_index == 8:
                        slide_placeholder_index = 2
                    else:
                        slide_placeholder_index = 1

                last_slide_layout_index = slide_layout_index
                continue

            elif line.startswith('#Header:'):
                header = line.replace('#Header:', '').strip()
                
                continue

            elif line.startswith('#Content:'):
                content = line.replace('#Content:', '').strip()
                next_line = f.readline().strip()
                while next_line and not next_line.startswith('#'):
                    content += '\n' + next_line
                    next_line = f.readline().strip()
                continue
                

    prs.save(f'GeneratedPresentations/{ppt_name}.pptx')
    
